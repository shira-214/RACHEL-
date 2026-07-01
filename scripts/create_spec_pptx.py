# -*- coding: utf-8 -*-
"""Generate PowerPoint from system specification document."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT = r"C:\Users\1\Downloads\פרוייקט\מסמך_אפיון_המשכיר.pptx"
FONT = "Segoe UI"
TITLE_COLOR = RGBColor(0, 120, 215)
ACCENT = RGBColor(51, 51, 51)


def set_rtl_paragraph(paragraph, text, size=18, bold=False, color=ACCENT):
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_rtl_paragraph(p, title, size=40, bold=True, color=TITLE_COLOR)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        set_rtl_paragraph(box2.text_frame.paragraphs[0], subtitle, size=20)


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    set_rtl_paragraph(box.text_frame.paragraphs[0], title, size=36, bold=True, color=RGBColor(255, 255, 255))


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    set_rtl_paragraph(title_box.text_frame.paragraphs[0], title, size=28, bold=True, color=TITLE_COLOR)

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item[1] if isinstance(item, tuple) else 0
        text = item[0] if isinstance(item, tuple) else item
        set_rtl_paragraph(p, "• " + text, size=16 if p.level == 0 else 14)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, 'מערכת "המשכיר"', "פרויקט גמר | אפיון מערכת | יולי 2026")

    add_bullet_slide(prs, "תקציר מנהלים", [
        "מערכת Desktop לניהול והזמנת דירות לשכירות קצרה",
        "שלושה תפקידים: שוכר, משכיר (בעל דירה), מנהל מערכת",
        "ארכיטקטורה: WPF Client + WCF Server + BL + DAL + EF6 + LocalDB",
    ])

    add_bullet_slide(prs, "מטרות המערכת", [
        "חיפוש והזמנה — שוכר מוצא דירה ומזמין לפי תאריך ומחיר",
        "ניהול נכסים — משכיר מפרסם דירות ומעלה תמונות",
        "ניהול מערכת — מנהל שולט בכל הנתונים (CRUD מלא)",
        "הפרדת תפקידים — כל משתמש רואה רק מה שמותר לו",
    ])

    add_section_slide(prs, "ארכיטקטורה טכנית")

    add_bullet_slide(prs, "שכבות המערכת", [
        "Presentation — WpfRentingApartementRacheli (מסכים, ניווט, Session)",
        "Service — Server / IService1 (חוזה WCF, CRUD)",
        "Host — Host.exe (הרצת שירות, Seed Data, כניסת מנהל)",
        "Business Logic — BL (המרות DTO, חוקי עסק, תמונות)",
        "Data Access — DAL (Entity Framework 6, Include)",
        "DTO — אובייקטי העברה בין שכבות",
        "מסד נתונים — RentingAppartmentRacheli.mdf (LocalDB)",
        "תמונות — תיקיית Pictures (קבצים + שם ב-DB)",
    ])

    add_bullet_slide(prs, "זרימת נתונים", [
        "WPF Client  ←→  Host.exe (WCF)  ←→  Service1",
        "Service1  →  BL  →  DAL  →  LocalDB",
        "תמונות: Client שולח bytes → Server שומר ב-Pictures → DB שומר שם קובץ",
        "כתובת שרת: localhost:8733",
    ])

    add_section_slide(prs, "מודל נתונים")

    add_bullet_slide(prs, "10 ישויות במערכת", [
        "Areas — אזורים גיאוגרפיים",
        "Cities — ערים",
        "StreetsNames — רחובות",
        "AraesCitiesStreet — קישור אזור-עיר-רחוב",
        "Apartments — דירות לשכירות",
        "Extras — תוספות (WiFi, חניה, מיזוג…)",
        "ExtrasApartements — קישור תוספת לדירה",
        "Images — תמונות לדירה",
        "Hirers — שוכרים (ת.ז.)",
        "Rentings — הזמנות / השכרות",
    ])

    add_bullet_slide(prs, "קשרים עיקריים", [
        "Apartment ← City, Street, Images, Extras, Rentings",
        "Renting ← Hirer + Apartment",
        "AraesCitiesStreet — לסינון גיאוגרפי בחיפוש",
        "שדות מפתח בדירה: מחיר, מיטות, סטטוס, בעלים, טלפון",
    ])

    add_section_slide(prs, "תפקידים והרשאות")

    add_bullet_slide(prs, "שוכר (Hirer)", [
        "כניסה: תעודת זהות (5–9 ספרות + בדיקת ספרת ביקורת)",
        "רישום: שם בעברית + טלפון נייד",
        "יכולות: חיפוש, סינון, הזמנה, צפייה ב\"ההזמנות שלי\"",
        "מסכים: pageLogin → PSearch → pageRenting → MyRentingsPage",
        "בדיקה: ת.ז. 123456782",
    ])

    add_bullet_slide(prs, "משכיר (Owner)", [
        "כניסה: שם בעלים + טלפון (NameOwner + PhoneOwner)",
        "רישום: הוספת דירה ראשונה",
        "יכולות: הוספה/עריכה של הדירות שלו, תמונות, צפייה בהזמנות",
        "מסכים: OwnerLoginPage → OwnerDashboardPage → PAdd / Image",
        "בדיקה: דוד כהן / 0501234567",
    ])

    add_bullet_slide(prs, "מנהל (Manager)", [
        "כניסה: admin / 1234 (Host App.config)",
        "יכולות: CRUD מלא על כל 10 הישויות",
        "מסכים: ManagerLoginPage → ManagerPage → All* / AddManager*",
        "מחיקה עם תלויות — הודעת שגיאה בעברית",
    ])

    add_section_slide(prs, "זרימות משתמש")

    add_bullet_slide(prs, "UC-01: חיפוש והזמנה (שוכר)", [
        "1. התחברות בת.ז.",
        "2. סינון: אזור, עיר, רחוב, מחיר, מיטות, תאריך, תוספת",
        "3. צפייה בכרטיסי דירות + תמונה",
        "4. הזמנה: מיטות, תאריך, פרטי אשראי",
        "5. בדיקת זמינות + קיבולת + ולידציה",
        "6. שמירת Renting + חזרה לחיפוש",
    ])

    add_bullet_slide(prs, "UC-02 / UC-03: משכיר ומנהל", [
        "משכיר: לוח בקרה → הוספת/עריכת דירה → העלאת תמונה (בחר + שמור)",
        "מנהל: כניסה admin → תפריט 10 ישויות → רשימה/הוספה/עדכון/מחיקה",
        "סדר מחיקה מומלץ: הזמנות → תמונות → תוספות → דירה",
    ])

    add_section_slide(prs, "כללי עסק")

    add_bullet_slide(prs, "Business Rules", [
        "BR-01: SumPayment = MinimumPrice + (ExtraForBed × SumBeds)",
        "BR-02: דירה תפוסה אם יש הזמנה באותו יום",
        "BR-03: בחיפוש — רק דירות עם Status = true",
        "BR-04: תמונה בכרטיס — פעילה (Stataus = true)",
        "BR-05: אשראי — ולידציה בלבד, ללא חיוב אמיתי",
        "BR-06: תמונה — קובץ ב-Pictures, DB שומר שם קובץ",
    ])

    add_bullet_slide(prs, "ולידציות", [
        "IsId — ת.ז. שוכר",
        "isHebrewRule — שם בעברית",
        "IsCellPhone — טלפון נייד",
        "כרטיס אשראי — 16 ספרות, CVV 3 ספרות, תוקף MM/YY",
    ])

    add_section_slide(prs, "מצב נוכחי")

    add_bullet_slide(prs, "מה עובד ✓", [
        "3 תפקידים עם הפרדת הרשאות",
        "CRUD מלא למנהל על 10 ישויות",
        "חיפוש עם סינון מתקדם",
        "הזמנה עם ולידציית אשראי",
        "העלאת תמונות (DB + Pictures)",
        "עברית + RTL, תיקוני bugs",
    ])

    add_bullet_slide(prs, "מגבלות ידועות", [
        "Session בזיכרון — אין persistence",
        "מנהל — סיסמה ב-config בלבד",
        "תשלום מדומה — אין סליקה אמיתית",
        "LocalDB — DB מקומי לכל מחשב",
        "WCF / EF6 — טכנולוגיה legacy",
    ])

    add_section_slide(prs, "הצעות לשיפור")

    add_bullet_slide(prs, "שיפורים מיידיים (לפני הגשה)", [
        "ComboBox דירה — הצגה ברורה (עיר, רחוב, טלפון)",
        "הודעות UX — בחר תמונה vs שמור תמונה",
        "Thumbnail בתצוגת תמונות למנהל",
        "בדיקת Host — הודעה אם השרת לא רץ",
        "רענון אוטומטי אחרי שמירה",
    ])

    add_bullet_slide(prs, "שיפורים בינוניים", [
        "גלריית תמונות לדירה",
        "סינון תוספות מרובות",
        "טווח תאריכים (check-in / check-out)",
        "דוחות מנהל — הכנסות, פופולריות",
        "Unit Tests ל-BL",
        "Logging לשגיאות שרת",
    ])

    add_bullet_slide(prs, "שדרוג ארוך טווח", [
        "ASP.NET Core Web API במקום WCF",
        "אפליקציית Web / Mobile (React, Blazor, MAUI)",
        "SQL Server / Azure SQL — DB מרכזי",
        "Azure Blob Storage — תמונות בענן",
        "JWT Authentication",
        "סליקה אמיתית (Stripe / Tranzila)",
        "התראות SMS/Email, דירוגים, מפת Google",
    ])

    add_bullet_slide(prs, "SWOT", [
        ("חוזקות: ארכיטקטורה שכבתית, 3 תפקידים, CRUD, ולידציות, עברית", 0),
        ("חולשות: Session מקומי, WCF ישן, DB מקומי, UX בסיסי", 0),
        ("הזדמנויות: Web, סליקה, אפליקציית נייד", 0),
        ("איומים: אובדן נתונים, קושי scale, תלות ב-Host", 0),
    ])

    add_bullet_slide(prs, "הרצה ותשתיות", [
        "1. Host.exe (חובה לפני WPF)",
        "2. WpfRentingApartementRacheli.exe",
        "DB: RentingAppartmentRacheli.mdf",
        "תמונות: Pictures\\",
        "Git: github.com/shira-214/RACHEL-",
    ])

    add_title_slide(prs, "תודה!", "מערכת \"המשכיר\" — פרויקט גמר 2026")

    prs.save(OUTPUT)
    print("Created:", OUTPUT)


if __name__ == "__main__":
    main()
